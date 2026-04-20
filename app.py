"""
pptx-service — AIWA Catalogue Generator
Pulls AIWA.pptx template from MinIO, populates product data, pushes result back.

POST /generate
{
  "brand": "AIWA",
  "category": "TWS Earphones",
  "products": [
    {
      "sku": "AI1001",
      "name": "CONNECT PRO TWS EARPHONES WITH CHARGING CASE",
      "description": "20 hours total playtime, IPX4 sweatproof, touch control",
      "image_url": "https://images.salsify.com/image/upload/s--xxx--/gacg3eq65mubj2hmt0w7.jpg",
      "upcs": [
        {"color": "BLK", "upc": "021331055405"},
        {"color": "WHT", "upc": "021331061447"}
      ],
      "case_pack": 12
    }
  ],
  "output_filename": "AIWA_TWS_Earphones_20260325.pptx"
}

Returns: { "url": "https://minio.asasson.xyz/..." }
"""

import os
import io
import html
import copy
import logging
import requests
from flask import Flask, request, jsonify
from pptx import Presentation
from lxml import etree
from minio import Minio
from minio.error import S3Error
from datetime import timedelta

# ── Config ────────────────────────────────────────────────────────────────────

MINIO_ENDPOINT   = os.environ.get("MINIO_ENDPOINT",   "minio:9000")
MINIO_ACCESS_KEY = os.environ.get("MINIO_ACCESS_KEY", "YOUR_ACCESS_KEY")
MINIO_SECRET_KEY = os.environ.get("MINIO_SECRET_KEY", "YOUR_SECRET_KEY")
MINIO_SECURE     = os.environ.get("MINIO_SECURE",     "false").lower() == "true"
MINIO_BUCKET     = os.environ.get("MINIO_BUCKET",     "sakar-catalogues")
TEMPLATE_PATH    = os.environ.get("TEMPLATE_PATH",    "AIWA.pptx")
OUTPUT_PREFIX    = "output/"
PUBLIC_MINIO_URL = os.environ.get("PUBLIC_MINIO_URL", "https://minio.asasson.xyz")

# ── App ───────────────────────────────────────────────────────────────────────

logging.basicConfig(level=logging.INFO)
log = logging.getLogger(__name__)
app = Flask(__name__)

# ── MinIO client ──────────────────────────────────────────────────────────────

minio_client = Minio(
    MINIO_ENDPOINT,
    access_key=MINIO_ACCESS_KEY,
    secret_key=MINIO_SECRET_KEY,
    secure=MINIO_SECURE,
)


def get_template() -> Presentation:
    """Download template from MinIO and return as Presentation object."""
    log.info(f"Fetching template: {TEMPLATE_PATH}")
    response = minio_client.get_object(MINIO_BUCKET, TEMPLATE_PATH)
    data = response.read()
    response.close()
    return Presentation(io.BytesIO(data))


def upload_result(prs: Presentation, filename: str) -> str:
    """Save Presentation to MinIO, return presigned URL valid 7 days."""
    buf = io.BytesIO()
    prs.save(buf)
    size = buf.tell()
    buf.seek(0)
    object_name = OUTPUT_PREFIX + filename
    minio_client.put_object(
        MINIO_BUCKET,
        object_name,
        buf,
        length=size,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    url = minio_client.presigned_get_object(
        MINIO_BUCKET, object_name, expires=timedelta(days=7)
    )
    # Replace internal endpoint with public URL
    url = url.replace(f"http://{MINIO_ENDPOINT}", PUBLIC_MINIO_URL)
    url = url.replace(f"https://{MINIO_ENDPOINT}", PUBLIC_MINIO_URL)
    log.info(f"Uploaded: {object_name}")
    return url


# ── Image helpers ─────────────────────────────────────────────────────────────

def download_image(url: str) -> bytes:
    """Download full-size image and compress to reasonable size."""
    import PIL.Image as PILImage

    # Remove thumbnail transform to get full size
    url = url.replace('/t_salsify_image_40/', '/')

    r = requests.get(url, timeout=30)
    r.raise_for_status()

    # Resize to max 600px wide and compress
    img = PILImage.open(io.BytesIO(r.content))
    w, h = img.size
    if w > 600:
        new_h = int(h * 600 / w)
        img = img.resize((600, new_h), PILImage.LANCZOS)

    # Convert RGBA to RGB if needed
    if img.mode in ('RGBA', 'P'):
        img = img.convert('RGB')

    buf = io.BytesIO()
    img.save(buf, format='JPEG', quality=82)
    return buf.getvalue()


def replace_picture(slide, shape, image_bytes: bytes):
    """
    Replace image and resize shape to square to avoid distorting square
    product images from Salsify. Shape is centered in its original vertical space.
    """
    import PIL.Image as PILImage
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.parts.image import ImagePart

    img = PILImage.open(io.BytesIO(image_bytes))
    fmt = (img.format or "JPEG").upper()
    ext_map = {"JPEG": "jpeg", "PNG": "png", "GIF": "gif"}
    ext = ext_map.get(fmt, "jpeg")
    mime_map = {"jpeg": "image/jpeg", "png": "image/png", "gif": "image/gif"}
    mime = mime_map.get(ext, "image/jpeg")

    package = slide.part.package
    partname = package.next_image_partname(ext)
    image_part = ImagePart(partname, mime, package, image_bytes)
    new_rId = slide.part.relate_to(image_part, RT.IMAGE)

    blip = shape._element.find(
        './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
    )
    if blip is not None:
        blip.set(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed',
            new_rId
        )

    # FIX: resize shape to square using width as the side length,
    # vertically centered within the original shape's space
    original_top = shape.top
    original_height = shape.height
    square_size = shape.width
    shape.height = square_size
    shape.top = original_top + (original_height - square_size) // 2


# ── Slide cloning ─────────────────────────────────────────────────────────────

def clone_slide(prs: Presentation, slide_index: int):
    """
    Clone an existing slide preserving all shapes, backgrounds and images.
    Appends the clone to the end of the presentation. Returns the new slide.
    """
    template = prs.slides[slide_index]
    slide_layout = template.slide_layout

    # Add new blank slide with same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Build rId mapping BEFORE copying shapes so blip refs can be fixed
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rId_map = {}
    for rId, rel in template.part.rels.items():
        if 'image' in rel.reltype:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    # Copy all shapes from template spTree
    template_spTree = template.shapes._spTree
    new_spTree = new_slide.shapes._spTree

    # Remove default placeholder shapes (keep nvGrpSpPr and grpSpPr at index 0,1)
    for child in list(new_spTree)[2:]:
        new_spTree.remove(child)

    # Deep copy shapes and fix blip rId references using the map
    for child in list(template_spTree)[2:]:
        new_child = copy.deepcopy(child)
        for blip in new_child.findall(f'.//{{{NS_A}}}blip'):
            old_rId = blip.get(f'{{{NS_R}}}embed')
            if old_rId and old_rId in rId_map:
                blip.set(f'{{{NS_R}}}embed', rId_map[old_rId])
        new_spTree.append(new_child)

    return new_slide


# ── XML cell builders ─────────────────────────────────────────────────────────

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _run_xml(text: str, bold: bool = False, sz: int = 1500) -> str:
    b = "1" if bold else "0"
    # FIX: unescape HTML entities from Salsify before re-escaping for XML
    # prevents double-encoding e.g. &amp; -> &amp;amp;
    text = html.unescape(text)
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return f"""<a:r xmlns:a="{NS}">
        <a:rPr lang="en-US" sz="{sz}" b="{b}" dirty="0">
          <a:solidFill><a:srgbClr val="000000"/></a:solidFill>
          <a:latin typeface="Calibri" panose="020F0502020204030204" pitchFamily="34" charset="0"/>
        </a:rPr>
        <a:t>{text}</a:t>
      </a:r>"""


def set_cell_text_simple(cell, text: str, bold: bool = False):
    """Replace cell content with a single paragraph, single run."""
    tc = cell._tc
    txBody = tc.find(f"{{{NS}}}txBody")
    for p in txBody.findall(f"{{{NS}}}p"):
        txBody.remove(p)
    p_xml = f"""<a:p xmlns:a="{NS}">
      <a:pPr algn="l">
        <a:lnSpc><a:spcPct val="100000"/></a:lnSpc>
      </a:pPr>
      {_run_xml(text, bold=bold)}
    </a:p>"""
    txBody.append(etree.fromstring(p_xml))


def set_cell_text_multiline(cell, text: str):
    """
    Split description on newlines, one paragraph per line.
    FIX: prevents long descriptions from rendering as one blob.
    """
    tc = cell._tc
    txBody = tc.find(f"{{{NS}}}txBody")
    for p in txBody.findall(f"{{{NS}}}p"):
        txBody.remove(p)
    lines = text.split("\n") if text else [""]
    for line in lines:
        line = line.strip()
        p_xml = f"""<a:p xmlns:a="{NS}">
          <a:pPr algn="l"><a:lnSpc><a:spcPct val="100000"/></a:lnSpc></a:pPr>
          {_run_xml(line, sz=1500)}
        </a:p>"""
        txBody.append(etree.fromstring(p_xml))


def set_cell_upc(cell, upcs: list):
    """Build UPC cell with one paragraph per color/upc pair."""
    tc = cell._tc
    txBody = tc.find(f"{{{NS}}}txBody")
    for p in txBody.findall(f"{{{NS}}}p"):
        txBody.remove(p)

    if not upcs:
        p_xml = f'<a:p xmlns:a="{NS}"><a:pPr algn="l"/></a:p>'
        txBody.append(etree.fromstring(p_xml))
        return

    for item in upcs:
        color = item.get("color", "")
        upc   = item.get("upc", "")
        line  = f"{color} UPC {upc}" if color else f"UPC {upc}"
        p_xml = f"""<a:p xmlns:a="{NS}">
          <a:pPr><a:lnSpc><a:spcPct val="100000"/></a:lnSpc></a:pPr>
          {_run_xml(line)}
        </a:p>"""
        txBody.append(etree.fromstring(p_xml))


# ── Category text box ─────────────────────────────────────────────────────────

def set_category_text(slide, category: str):
    """Find the category text box (top right) and update its text."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 17:  # TEXT_BOX
            if shape.text_frame.text.strip():
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = category
                return
    log.warning("Category text box not found")


# ── Populate one line-sheet slide ─────────────────────────────────────────────

def populate_line_sheet(slide, category: str, products: list):
    """Fill a cloned line-sheet slide with product data."""

    # 1. Update category name
    set_category_text(slide, category)

    # 2. Collect picture shapes sorted left to right
    pictures = sorted(
        [s for s in slide.shapes if s.shape_type == 13],
        key=lambda s: s.left
    )

    # 3. Find table
    table_shape = next((s for s in slide.shapes if s.shape_type == 19), None)
    if table_shape is None:
        log.error("No table found on slide")
        return

    table = table_shape.table
    num_cols = len(table.columns)

    # FIX: cap to actual picture count — prevents addressing non-picture shapes
    num_products = min(len(products), len(pictures), num_cols)

    # 4. Replace images
    for i in range(num_products):
        product = products[i]
        image_url = product.get("image_url", "")
        if i < len(pictures) and image_url:
            try:
                img_bytes = download_image(image_url)
                replace_picture(slide, pictures[i], img_bytes)
                log.info(f"  Image {i} replaced: {product['sku']}")
            except Exception as e:
                log.warning(f"  Image {i} failed ({product['sku']}): {e}")

    # 5. Populate table cells
    for i in range(num_products):
        p = products[i]
        set_cell_text_simple(table.rows[0].cells[i], p.get("sku", ""), bold=True)
        set_cell_text_simple(table.rows[1].cells[i], p.get("name", "").upper())
        # FIX: use multiline builder for description
        set_cell_text_multiline(table.rows[2].cells[i], p.get("description", ""))
        set_cell_upc(table.rows[3].cells[i], p.get("upcs", []))
        cp = p.get("case_pack", "")
        set_cell_text_simple(table.rows[4].cells[i], f"CP {cp}" if cp else "")

    # 6. Clear unused columns
    for i in range(num_products, num_cols):
        for row_idx in range(5):
            set_cell_text_simple(table.rows[row_idx].cells[i], "")
        if i < len(pictures):
            sp = pictures[i]
            sp._element.getparent().remove(sp._element)


# ── Build full catalogue ──────────────────────────────────────────────────────

def chunk_products(products: list, size: int = 6) -> list:
    return [products[i:i+size] for i in range(0, len(products), size)]


def build_catalogue(prs: Presentation, category: str, products: list) -> Presentation:
    """
    Keep slides 1+2, replace slides 3+ with generated line sheets.
    Uses slide 3 as the 6-up line sheet template.
    """
    TEMPLATE_SLIDE_IDX = 2  # slide 3 = 6-up line sheet

    chunks = chunk_products(products)

    # FIX: capture original slide count BEFORE cloning anything
    # Previously used len(prs.slides) after cloning, which included the new
    # slides and caused the deletion range to be wrong, overwriting template slides
    original_slide_count = len(prs.slides)

    # Clone template slide once per chunk
    new_slides = []
    for chunk in chunks:
        new_slide = clone_slide(prs, TEMPLATE_SLIDE_IDX)
        new_slides.append((new_slide, chunk))

    # Remove original product slides (index 2 onwards), using the
    # pre-clone count so we don't touch the newly appended slides
    slides_to_keep = 2  # cover + brand history
    sldIdLst = prs.slides._sldIdLst
    all_sldIds = list(sldIdLst)

    for sldId in all_sldIds[slides_to_keep:original_slide_count]:
        sldIdLst.remove(sldId)

    # Populate each new slide with product data
    for new_slide, chunk in new_slides:
        populate_line_sheet(new_slide, category, chunk)

    return prs


# ── Flask routes ──────────────────────────────────────────────────────────────

@app.route("/generate", methods=["POST"])
def generate():
    body = request.get_json(force=True)
    brand    = body.get("brand", "AIWA")
    category = body.get("category", "Products")
    products = body.get("products", [])
    filename = body.get("output_filename", f"{brand}_{category.replace(' ', '_')}.pptx")

    if not products:
        return jsonify({"error": "No products provided"}), 400

    log.info(f"Generating: {brand} / {category} / {len(products)} products")

    try:
        prs = get_template()
        prs = build_catalogue(prs, category, products)
        url = upload_result(prs, filename)
        return jsonify({"url": url, "filename": filename, "product_count": len(products)})
    except S3Error as e:
        log.error(f"MinIO error: {e}")
        return jsonify({"error": f"Storage error: {str(e)}"}), 500
    except Exception as e:
        log.exception("Generation failed")
        return jsonify({"error": str(e)}), 500


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5001, debug=False)

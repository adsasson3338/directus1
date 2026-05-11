"""
pptx-service — AIWA Catalogue Generator
Pulls AIWA.pptx template from MinIO, populates product data, pushes result back.

POST /generate
{
  "brand": "AIWA",
  "category": "TWS Earphones",
  "products": [
    {
      "sku": "AI1001",
      "name": "CONNECT PRO TWS EARPHONES WITH CHARGING CASE",
      "description": "20 hours total playtime, IPX4 sweatproof, touch control",
      "image_url": "https://images.salsify.com/image/upload/s--xxx--/gacg3eq65mubj2hmt0w7.jpg",
      "upcs": [
        {"color": "BLK", "upc": "021331055405"},
        {"color": "WHT", "upc": "021331061447"}
      ],
      "case_pack": 12
    }
  ],
  "output_filename": "AIWA_TWS_Earphones_20260325.pptx"
}

Returns: { "url": "https://minio.asasson.xyz/..." }
"""

import os
import io
import html
import copy
import logging
import requests
from flask import Flask, request, jsonify
from pptx import Presentation
from lxml import etree
from minio import Minio
from minio.error import S3Error
from datetime import timedelta

# ── Config ────────────────────────────────────────────────────────────────────

MINIO_ENDPOINT   = os.environ.get("MINIO_ENDPOINT",   "minio:9000")
MINIO_ACCESS_KEY = os.environ.get("MINIO_ACCESS_KEY", "YOUR_ACCESS_KEY")
MINIO_SECRET_KEY = os.environ.get("MINIO_SECRET_KEY", "YOUR_SECRET_KEY")
MINIO_SECURE     = os.environ.get("MINIO_SECURE",     "false").lower() == "true"
MINIO_BUCKET     = os.environ.get("MINIO_BUCKET",     "sakar-catalogues")
TEMPLATE_PATH    = os.environ.get("TEMPLATE_PATH",    "AIWA.pptx")
OUTPUT_PREFIX    = "output/"
PUBLIC_MINIO_URL = os.environ.get("PUBLIC_MINIO_URL", "https://minio.asasson.xyz")

# ── App ───────────────────────────────────────────────────────────────────────

logging.basicConfig(level=logging.INFO)
log = logging.getLogger(__name__)
app = Flask(__name__)

# ── MinIO client ──────────────────────────────────────────────────────────────

minio_client = Minio(
    MINIO_ENDPOINT,
    access_key=MINIO_ACCESS_KEY,
    secret_key=MINIO_SECRET_KEY,
    secure=MINIO_SECURE,
)


def get_template() -> Presentation:
    """Download template from MinIO and return as Presentation object."""
    log.info(f"Fetching template: {TEMPLATE_PATH}")
    response = minio_client.get_object(MINIO_BUCKET, TEMPLATE_PATH)
    data = response.read()
    response.close()
    return Presentation(io.BytesIO(data))


def upload_result(prs: Presentation, filename: str) -> str:
    """Save Presentation to MinIO, return presigned URL valid 7 days."""
    buf = io.BytesIO()
    prs.save(buf)
    size = buf.tell()
    buf.seek(0)
    object_name = OUTPUT_PREFIX + filename
    minio_client.put_object(
        MINIO_BUCKET,
        object_name,
        buf,
        length=size,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    url = minio_client.presigned_get_object(
        MINIO_BUCKET, object_name, expires=timedelta(days=7)
    )
    # Replace internal endpoint with public URL
    url = url.replace(f"http://{MINIO_ENDPOINT}", PUBLIC_MINIO_URL)
    url = url.replace(f"https://{MINIO_ENDPOINT}", PUBLIC_MINIO_URL)
    log.info(f"Uploaded: {object_name}")
    return url


# ── Image helpers ─────────────────────────────────────────────────────────────

def download_image(url: str) -> bytes:
    """Download image, cap longest side to 600px, always return JPEG bytes."""
    import PIL.Image as PILImage

    # Remove thumbnail transform to get full size
    url = url.replace('/t_salsify_image_40/', '/')

    r = requests.get(url, timeout=30)
    r.raise_for_status()

    img = PILImage.open(io.BytesIO(r.content))
    w, h = img.size

    # Cap longest side at 600px, preserving aspect ratio
    MAX = 600
    if w > MAX or h > MAX:
        scale = MAX / max(w, h)
        img = img.resize((int(w * scale), int(h * scale)), PILImage.LANCZOS)

    # Convert to RGB (handles RGBA, P, LA, etc.)
    if img.mode in ('RGBA', 'P', 'LA'):
        img = img.convert('RGBA')
        background = PILImage.new('RGB', img.size, (255, 255, 255))
        background.paste(img, mask=img.getchannel('A'))
        img = background
    elif img.mode != 'RGB':
        img = img.convert('RGB')

    buf = io.BytesIO()
    img.save(buf, format='JPEG', quality=82)
    return buf.getvalue()


def replace_picture(slide, shape, image_bytes: bytes):
    """
    Swap the image in a picture shape.
    - Preserves aspect ratio (no distortion)
    - Fits within the original shape bounding box, centered
    - Removes any border inherited from the template
    """
    import PIL.Image as PILImage
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.parts.image import ImagePart

    NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    # ── 1. Register new image with the slide ──────────────────────────────────
    package  = slide.part.package
    partname = package.next_image_partname("jpeg")
    image_part = ImagePart(partname, "image/jpeg", package, image_bytes)
    new_rId = slide.part.relate_to(image_part, RT.IMAGE)

    # ── 2. Point the blip at the new image ────────────────────────────────────
    blip = shape._element.find(f'.//{{{NS_A}}}blip')
    if blip is not None:
        blip.set(f'{{{NS_R}}}embed', new_rId)

    # ── 3. Remove border (template has a solid tx1 line on every picture) ─────
    spPr = shape._element.find(f'{{{NS_P}}}spPr')
    if spPr is not None:
        for ln in spPr.findall(f'{{{NS_A}}}ln'):
            spPr.remove(ln)
        spPr.append(etree.fromstring(f'<a:ln xmlns:a="{NS_A}"><a:noFill/></a:ln>'))

    # ── 4. Resize shape to match image aspect ratio, centered in original box ─
    img = PILImage.open(io.BytesIO(image_bytes))
    img_w, img_h = img.size

    orig_left   = shape.left
    orig_top    = shape.top
    orig_width  = shape.width
    orig_height = shape.height

    img_aspect = img_w / img_h
    box_aspect = orig_width / orig_height

    if img_aspect > box_aspect:
        # Wider than box — fit to width
        new_w = orig_width
        new_h = int(orig_width / img_aspect)
    else:
        # Taller than box — fit to height
        new_h = orig_height
        new_w = int(orig_height * img_aspect)

    shape.width  = new_w
    shape.height = new_h
    shape.left   = orig_left + (orig_width  - new_w) // 2
    shape.top    = orig_top  + (orig_height - new_h) // 2


# ── Slide cloning ─────────────────────────────────────────────────────────────

def clone_slide(prs: Presentation, slide_index: int):
    """
    Clone an existing slide preserving all shapes, backgrounds and images.
    Appends the clone to the end of the presentation. Returns the new slide.
    """
    template     = prs.slides[slide_index]
    slide_layout = template.slide_layout
    new_slide    = prs.slides.add_slide(slide_layout)

    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Build rId mapping BEFORE copying shapes so blip refs can be fixed
    rId_map = {}
    for rId, rel in template.part.rels.items():
        if 'image' in rel.reltype:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    # Copy all shapes from template, removing default placeholders first
    template_spTree = template.shapes._spTree
    new_spTree      = new_slide.shapes._spTree

    for child in list(new_spTree)[2:]:
        new_spTree.remove(child)

    for child in list(template_spTree)[2:]:
        new_child = copy.deepcopy(child)
        for blip in new_child.findall(f'.//{{{NS_A}}}blip'):
            old_rId = blip.get(f'{{{NS_R}}}embed')
            if old_rId and old_rId in rId_map:
                blip.set(f'{{{NS_R}}}embed', rId_map[old_rId])
        new_spTree.append(new_child)

    return new_slide


# ── XML cell builders ─────────────────────────────────────────────────────────

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _run_xml(text: str, bold: bool = False, sz: int = 1500) -> str:
    b = "1" if bold else "0"
    # Unescape HTML entities from Salsify before re-escaping for XML
    text = html.unescape(text)
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return f"""<a:r xmlns:a="{NS}">
        <a:rPr lang="en-US" sz="{sz}" b="{b}" dirty="0">
          <a:solidFill><a:srgbClr val="000000"/></a:solidFill>
          <a:latin typeface="Calibri" panose="020F0502020204030204" pitchFamily="34" charset="0"/>
        </a:rPr>
        <a:t>{text}</a:t>
      </a:r>"""


def set_cell_text_simple(cell, text: str, bold: bool = False):
    """Replace cell content with a single paragraph, single run."""
    tc     = cell._tc
    txBody = tc.find(f"{{{NS}}}txBody")
    for p in txBody.findall(f"{{{NS}}}p"):
        txBody.remove(p)
    p_xml = f"""<a:p xmlns:a="{NS}">
      <a:pPr algn="l"><a:lnSpc><a:spcPct val="100000"/></a:lnSpc></a:pPr>
      {_run_xml(text, bold=bold)}
    </a:p>"""
    txBody.append(etree.fromstring(p_xml))


def set_cell_text_multiline(cell, text: str):
    """Split on newlines, one paragraph per line."""
    tc     = cell._tc
    txBody = tc.find(f"{{{NS}}}txBody")
    for p in txBody.findall(f"{{{NS}}}p"):
        txBody.remove(p)
    lines = text.split("\n") if text else [""]
    for line in lines:
        line  = line.strip()
        p_xml = f"""<a:p xmlns:a="{NS}">
          <a:pPr algn="l"><a:lnSpc><a:spcPct val="100000"/></a:lnSpc></a:pPr>
          {_run_xml(line, sz=1500)}
        </a:p>"""
        txBody.append(etree.fromstring(p_xml))


def set_cell_upc(cell, upcs: list):
    """Build UPC cell with one paragraph per color/upc pair."""
    tc     = cell._tc
    txBody = tc.find(f"{{{NS}}}txBody")
    for p in txBody.findall(f"{{{NS}}}p"):
        txBody.remove(p)

    if not upcs:
        txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:pPr algn="l"/></a:p>'))
        return

    for item in upcs:
        color = item.get("color", "")
        upc   = item.get("upc", "")
        line  = f"{color} UPC {upc}" if color else f"UPC {upc}"
        p_xml = f"""<a:p xmlns:a="{NS}">
          <a:pPr><a:lnSpc><a:spcPct val="100000"/></a:lnSpc></a:pPr>
          {_run_xml(line)}
        </a:p>"""
        txBody.append(etree.fromstring(p_xml))


# ── Category text box ─────────────────────────────────────────────────────────

def set_category_text(slide, category: str):
    """Find the category text box (top right) and update its text."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 17:  # TEXT_BOX
            if shape.text_frame.text.strip():
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = category
                return
    log.warning("Category text box not found")


# ── Populate one line-sheet slide ─────────────────────────────────────────────

def populate_line_sheet(slide, category: str, products: list):
    """
    Fill a cloned line-sheet slide with product data.
    If fewer than 6 products, unused columns and images are removed
    and the remaining content is centered on the slide.
    """

    # 1. Category label
    set_category_text(slide, category)

    # 2. Collect picture shapes left-to-right
    pictures = sorted(
        [s for s in slide.shapes if s.shape_type == 13],
        key=lambda s: s.left
    )

    # 3. Find table
    table_shape = next((s for s in slide.shapes if s.shape_type == 19), None)
    if table_shape is None:
        log.error("No table found on slide")
        return

    table        = table_shape.table
    num_cols     = len(table.columns)
    num_products = min(len(products), len(pictures), num_cols)
    total_slots  = len(pictures)

    # 4. If this is a partial page, trim unused columns and images
    is_partial = num_products < total_slots
    if is_partial:
        NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        tbl_xml = table._tbl

        # Remove unused picture shapes (right-to-left)
        for i in range(total_slots - 1, num_products - 1, -1):
            pictures[i]._element.getparent().remove(pictures[i]._element)

        # Remove extra gridCol entries from the table
        grid      = tbl_xml.find(f'{{{NS_A}}}tblGrid')
        grid_cols = grid.findall(f'{{{NS_A}}}gridCol')
        for col in grid_cols[num_products:]:
            grid.remove(col)

        # Remove extra cells from each row
        for row in tbl_xml.findall(f'{{{NS_A}}}tr'):
            cells = row.findall(f'{{{NS_A}}}tc')
            for cell in cells[num_products:]:
                row.remove(cell)

        # Redistribute column widths evenly across the original table width
        total_table_width = table_shape.width
        new_col_width     = total_table_width // num_products
        remainder         = total_table_width - (new_col_width * num_products)

        for i, gc in enumerate(grid.findall(f'{{{NS_A}}}gridCol')):
            w = new_col_width + (remainder if i == num_products - 1 else 0)
            gc.set('w', str(w))

        # Center the table horizontally on the slide
        slide_width       = slide.part.presentation.slide_width
        table_shape.left  = (slide_width - total_table_width) // 2

    # 5. Replace images — this resizes each shape to fit the image
    pictures_now = sorted(
        [s for s in slide.shapes if s.shape_type == 13],
        key=lambda s: s.left
    )
    for i in range(num_products):
        image_url = products[i].get("image_url", "")
        if image_url:
            try:
                img_bytes = download_image(image_url)
                replace_picture(slide, pictures_now[i], img_bytes)
                log.info(f"  Image {i} replaced: {products[i]['sku']}")
            except Exception as e:
                log.warning(f"  Image {i} failed ({products[i]['sku']}): {e}")

    # 6. Re-center images over their columns AFTER replace_picture has sized them
    #    (replace_picture changes shape dimensions, so we align after the fact)
    if is_partial:
        pictures_now = sorted(
            [s for s in slide.shapes if s.shape_type == 13],
            key=lambda s: s.left
        )
        for i, pic in enumerate(pictures_now):
            col_center = table_shape.left + (i * new_col_width) + (new_col_width // 2)
            pic.left   = col_center - (pic.width // 2)

    # 7. Fill table cells
    for i in range(num_products):
        p = products[i]
        set_cell_text_simple(table.rows[0].cells[i], p.get("sku", ""),           bold=True)
        set_cell_text_simple(table.rows[1].cells[i], p.get("name", "").upper())
        set_cell_text_multiline(table.rows[2].cells[i], p.get("description", ""))
        set_cell_upc(table.rows[3].cells[i],           p.get("upcs", []))
        cp = p.get("case_pack", "")
        set_cell_text_simple(table.rows[4].cells[i], f"CP {cp}" if cp else "")


# ── Build full catalogue ──────────────────────────────────────────────────────

def chunk_products(products: list, size: int = 6) -> list:
    return [products[i:i+size] for i in range(0, len(products), size)]


def build_catalogue(prs: Presentation, category: str, products: list) -> Presentation:
    """
    Keep slides 1+2, replace slides 3+ with generated line sheets.
    Uses slide 3 as the 6-up line sheet template.
    """
    TEMPLATE_SLIDE_IDX = 2  # slide 3 = 6-up line sheet

    chunks = chunk_products(products)

    # Capture original slide count BEFORE cloning so deletion range stays correct
    original_slide_count = len(prs.slides)

    # Clone template slide once per chunk
    new_slides = []
    for chunk in chunks:
        new_slide = clone_slide(prs, TEMPLATE_SLIDE_IDX)
        new_slides.append((new_slide, chunk))

    # Remove original product slides (keep cover + brand history)
    slides_to_keep = 2
    sldIdLst   = prs.slides._sldIdLst
    all_sldIds = list(sldIdLst)
    for sldId in all_sldIds[slides_to_keep:original_slide_count]:
        sldIdLst.remove(sldId)

    # Populate each cloned slide
    for new_slide, chunk in new_slides:
        populate_line_sheet(new_slide, category, chunk)

    return prs


# ── Flask routes ──────────────────────────────────────────────────────────────

@app.route("/generate", methods=["POST"])
def generate():
    body     = request.get_json(force=True)
    brand    = body.get("brand", "AIWA")
    category = body.get("category", "Products")
    products = body.get("products", [])
    filename = body.get("output_filename", f"{brand}_{category.replace(' ', '_')}.pptx")

    if not products:
        return jsonify({"error": "No products provided"}), 400

    log.info(f"Generating: {brand} / {category} / {len(products)} products")

    try:
        prs = get_template()
        prs = build_catalogue(prs, category, products)
        url = upload_result(prs, filename)
        return jsonify({"url": url, "filename": filename, "product_count": len(products)})
    except S3Error as e:
        log.error(f"MinIO error: {e}")
        return jsonify({"error": f"Storage error: {str(e)}"}), 500
    except Exception as e:
        log.exception("Generation failed")
        return jsonify({"error": str(e)}), 500


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5001, debug=False)
